#!/usr/bin/env python3
# SPDX-License-Identifier: CC-BY-SA-4.0
# (c) 2026 Mega Promoting SRL. Licensed under Creative Commons Attribution-ShareAlike 4.0 International.
"""generate-pptx.py — Build the MD-Chat Moldova Digital Summit deck.

Reads `MD-Chat-Moldova-Digital-Summit.md` from the same directory and
produces a branded 13-slide PPTX (1 title + 12 content).

Brand palette (matches /brand/colors.css dark theme):
  - Navy background   #1A2D4E
  - Teal accent       #2DD4BF
  - Snow text         #F8FAFC
  - Slate sub-text    #94A3B8

Requires:
  - Python 3.10+
  - python-pptx >= 1.0  (`pip install -r requirements.txt`)

CLI:
  python3 generate-pptx.py --output ./MD-Chat-pitch.pptx
  python3 generate-pptx.py            # default output: ./MD-Chat-pitch.pptx

Idempotent: re-running overwrites the output file.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path
from typing import List, Tuple

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError as exc:  # pragma: no cover
    sys.stderr.write(
        "ERROR: python-pptx not installed. Run: pip install -r requirements.txt\n"
    )
    raise SystemExit(1) from exc


# ---------------------------------------------------------------------------
# Brand constants (mirror /brand/colors.css :root.dark)
# ---------------------------------------------------------------------------

NAVY = RGBColor(0x1A, 0x2D, 0x4E)
NAVY_ELEVATED = RGBColor(0x22, 0x36, 0x5A)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
SNOW = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x94, 0xA3, 0xB8)
GOLD = RGBColor(0xFB, 0xBF, 0x24)

FONT_HEAD = "Inter"           # falls back to system sans-serif if absent
FONT_BODY = "Inter"
FONT_MONO = "JetBrains Mono"

# 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ---------------------------------------------------------------------------
# Markdown parser — light-weight, tailored for OUR file
# ---------------------------------------------------------------------------

SLIDE_HEADER_RE = re.compile(r"^##\s+Slide\s+(\d+)\s+[—–-]\s+(.+?)\s*$")


def parse_slides(md_path: Path) -> List[Tuple[str, List[str], str]]:
    """Return a list of (title, body_lines, speaker_notes) tuples.

    body_lines are stripped of leading bullet markers (`- `, `* `) but keep
    inline markdown for later cleanup. blockquotes (`>`) are treated as
    speaker notes and pulled aside.
    """
    lines = md_path.read_text(encoding="utf-8").splitlines()
    slides: List[Tuple[str, List[str], str]] = []

    cur_title = None
    cur_body: List[str] = []
    cur_notes: List[str] = []

    def flush():
        nonlocal cur_title, cur_body, cur_notes
        if cur_title is not None:
            slides.append((cur_title, cur_body, "\n".join(cur_notes).strip()))
        cur_title = None
        cur_body = []
        cur_notes = []

    for raw in lines:
        m = SLIDE_HEADER_RE.match(raw)
        if m:
            flush()
            cur_title = m.group(2).strip()
            continue

        if cur_title is None:
            continue

        if raw.strip().startswith("---"):
            flush()
            continue

        if raw.lstrip().startswith(">"):
            cur_notes.append(raw.lstrip()[1:].strip())
            continue

        cur_body.append(raw)

    flush()
    return slides


_BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
_CODE_RE = re.compile(r"`([^`]+)`")
_BULLET_RE = re.compile(r"^\s*[-*]\s+")


def clean_body_text(raw: str) -> str:
    """Strip markdown markers for display in pptx (formatting via runs is
    expensive; we keep text plain but tidy)."""
    s = raw
    s = _BULLET_RE.sub("", s)
    s = _BOLD_RE.sub(lambda m: m.group(1), s)
    s = _CODE_RE.sub(lambda m: m.group(1), s)
    s = s.strip()
    return s


def is_bullet(raw: str) -> bool:
    return bool(_BULLET_RE.match(raw))


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------


def set_slide_background(slide, color: RGBColor) -> None:
    """Solid fill the slide background with `color`."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, *, top: Inches, color: RGBColor = TEAL) -> None:
    """Teal accent bar — short horizontal rule under the title."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        top,
        Inches(1.2),
        Emu(40_000),  # ~0.04in thick
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = color


def add_textbox(
    slide,
    text: str,
    *,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
    font_name: str = FONT_HEAD,
    font_size: Pt = Pt(20),
    color: RGBColor = SNOW,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_body_bullets(
    slide,
    body_raw_lines: List[str],
    *,
    left: Inches = Inches(0.7),
    top: Inches = Inches(2.3),
    width: Inches = Inches(12.0),
    height: Inches = Inches(4.7),
):
    """Render body as bullet list. Skips empty lines, demotes
    non-bullet lines to indented paragraphs without bullets."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    first = True
    any_content = False
    for raw in body_raw_lines:
        stripped = raw.strip()
        if not stripped:
            continue
        # Skip fenced code marker lines.
        if stripped.startswith("```"):
            continue

        cleaned = clean_body_text(stripped)
        if not cleaned:
            continue

        any_content = True
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        if is_bullet(raw):
            run = p.add_run()
            run.text = "•  " + cleaned
            run.font.name = FONT_BODY
            run.font.size = Pt(20)
            run.font.color.rgb = SNOW
        elif stripped.startswith("#"):
            # demoted heading inside a slide (e.g. one-pager subsections)
            run = p.add_run()
            run.text = cleaned.lstrip("# ").strip()
            run.font.name = FONT_HEAD
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = TEAL
        else:
            run = p.add_run()
            run.text = cleaned
            run.font.name = FONT_BODY
            run.font.size = Pt(20)
            run.font.color.rgb = SNOW

        p.space_after = Pt(6)

    if not any_content:
        # Avoid stray empty textbox
        sp = tb._element
        sp.getparent().remove(sp)


def add_speaker_notes(slide, notes_text: str) -> None:
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_footer(slide, index: int, total: int = 12) -> None:
    """Small footer: deck title + slide N/total."""
    add_textbox(
        slide,
        "MD-Chat — Moldova Digital Summit 2026",
        left=Inches(0.7),
        top=Inches(7.05),
        width=Inches(8.0),
        height=Inches(0.35),
        font_name=FONT_BODY,
        font_size=Pt(10),
        color=SLATE,
    )
    add_textbox(
        slide,
        f"{index} / {total}",
        left=Inches(11.8),
        top=Inches(7.05),
        width=Inches(1.0),
        height=Inches(0.35),
        font_name=FONT_BODY,
        font_size=Pt(10),
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Title slide
# ---------------------------------------------------------------------------


def build_title_slide(prs: Presentation) -> None:
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, NAVY)

    # Big teal accent block on the left
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(0.4),
        SLIDE_HEIGHT,
    )
    block.line.fill.background()
    block.fill.solid()
    block.fill.fore_color.rgb = TEAL

    # Logo placeholder rectangle (1.5in square) — user can paste real logo
    logo = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.9),
        Inches(0.9),
        Inches(1.5),
        Inches(1.5),
    )
    logo.line.color.rgb = TEAL
    logo.line.width = Pt(2)
    logo.fill.solid()
    logo.fill.fore_color.rgb = NAVY_ELEVATED
    # Centered "MD" placeholder
    logo_tf = logo.text_frame
    logo_tf.margin_left = Emu(0)
    logo_tf.margin_right = Emu(0)
    p = logo_tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "MD"
    r.font.name = FONT_HEAD
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = TEAL

    # Main title
    add_textbox(
        slide,
        "MD-Chat",
        left=Inches(0.9),
        top=Inches(3.0),
        width=Inches(12.0),
        height=Inches(1.4),
        font_name=FONT_HEAD,
        font_size=Pt(72),
        color=SNOW,
        bold=True,
    )
    # Subtitle
    add_textbox(
        slide,
        "Mesager sovereign EU-grade, construit în Moldova.",
        left=Inches(0.9),
        top=Inches(4.2),
        width=Inches(12.0),
        height=Inches(0.8),
        font_name=FONT_BODY,
        font_size=Pt(28),
        color=TEAL,
    )
    add_textbox(
        slide,
        "Open source · AI confidențial nativ · EVO + MPass din Day 1",
        left=Inches(0.9),
        top=Inches(4.95),
        width=Inches(12.0),
        height=Inches(0.6),
        font_name=FONT_BODY,
        font_size=Pt(20),
        color=SLATE,
    )

    # Event line
    add_textbox(
        slide,
        "Moldova Digital Summit · 5–6 iunie 2026 · Chișinău",
        left=Inches(0.9),
        top=Inches(6.0),
        width=Inches(12.0),
        height=Inches(0.5),
        font_name=FONT_BODY,
        font_size=Pt(16),
        color=SNOW,
    )
    # Author line
    add_textbox(
        slide,
        "Oleg Chetrean · Mega Promoting SRL · oleg@megapromoting.com",
        left=Inches(0.9),
        top=Inches(6.45),
        width=Inches(12.0),
        height=Inches(0.5),
        font_name=FONT_MONO,
        font_size=Pt(13),
        color=SLATE,
    )

    add_speaker_notes(
        slide,
        "Title slide. Lasă afișat 10 secunde înainte de a începe. "
        "Auditoriul citește. Apoi treci la Slide 1 (Hook).",
    )


# ---------------------------------------------------------------------------
# Content slide builder
# ---------------------------------------------------------------------------


def build_content_slide(
    prs: Presentation, index: int, title: str, body: List[str], notes: str
) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_background(slide, NAVY)

    # Left side teal stripe (thin)
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(0.15),
        SLIDE_HEIGHT,
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = TEAL

    # Slide number marker (top-right)
    add_textbox(
        slide,
        f"Slide {index:02d}",
        left=Inches(11.2),
        top=Inches(0.45),
        width=Inches(1.6),
        height=Inches(0.4),
        font_name=FONT_MONO,
        font_size=Pt(12),
        color=SLATE,
        align=PP_ALIGN.RIGHT,
    )

    # Title
    add_textbox(
        slide,
        title,
        left=Inches(0.7),
        top=Inches(0.6),
        width=Inches(11.0),
        height=Inches(1.2),
        font_name=FONT_HEAD,
        font_size=Pt(36),
        color=SNOW,
        bold=True,
    )

    # Teal accent under title
    add_accent_bar(slide, top=Inches(1.9))

    # Body
    add_body_bullets(slide, body)

    # Footer
    add_footer(slide, index)

    # Speaker notes
    add_speaker_notes(slide, notes)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build_deck(md_path: Path, out_path: Path) -> Path:
    if not md_path.exists():
        raise FileNotFoundError(f"Source markdown not found: {md_path}")

    slides = parse_slides(md_path)
    if len(slides) < 12:
        sys.stderr.write(
            f"WARN: expected >=12 slides, parsed {len(slides)} from {md_path.name}\n"
        )

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_title_slide(prs)

    for i, (title, body, notes) in enumerate(slides[:12], start=1):
        build_content_slide(prs, i, title, body, notes)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main(argv: List[str] | None = None) -> int:
    here = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(
        description="Generate MD-Chat Moldova Digital Summit PPTX deck.",
    )
    parser.add_argument(
        "--source",
        type=Path,
        default=here / "MD-Chat-Moldova-Digital-Summit.md",
        help="Source markdown (default: ./MD-Chat-Moldova-Digital-Summit.md)",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=here / "MD-Chat-pitch.pptx",
        help="Output PPTX path (default: ./MD-Chat-pitch.pptx)",
    )
    args = parser.parse_args(argv)

    out = build_deck(args.source, args.output)
    print(f"OK  PPTX written: {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
